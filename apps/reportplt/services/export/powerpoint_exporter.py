 # apps/reportplt/services/export/powerpoint_exporter.py
import io
import uuid
from typing import Dict, Any, List, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from apps.reportplt.exceptions import ReportExportError

class PowerPointExporter:
    def __init__(self, config: Optional[Dict] = None):
        self.config = config or {}
        self.prs = Presentation()
        self._setup_slide_layouts()

    def _setup_slide_layouts(self):
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.content_slide_layout = self.prs.slide_layouts[1]
        self.title_only_slide_layout = self.prs.slide_layouts[5]

    def export(self, data: Dict[str, Any], report_name: str, output_path: Optional[str] = None) -> str:
        try:
            self._add_title_slide(report_name)
            self._add_executive_summary_slide(data)
            self._add_kpi_slides(data)
            self._add_chart_slides(data)
            self._add_table_slides(data)
            self._add_mission_status_slides(data)
            self._add_conclusion_slide()
            buffer = io.BytesIO()
            self.prs.save(buffer)
            buffer.seek(0)
            if output_path:
                with default_storage.open(output_path, 'wb') as f:
                    f.write(buffer.getvalue())
                return output_path
            file_name = f"reports/{uuid.uuid4()}.pptx"
            path = default_storage.save(file_name, ContentFile(buffer.getvalue()))
            return path
        except Exception as e:
            raise ReportExportError(f"PowerPoint export failed: {str(e)}")

    def _add_title_slide(self, report_name: str):
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = report_name
        subtitle.text = f"Performance Management Report\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}\nConfidential"

    def _add_executive_summary_slide(self, data: Dict):
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        title.text = "Executive Summary"
        content = slide.placeholders[1]
        summary = data.get('executive_summary', 'No summary available')
        content.text = summary

    def _add_kpi_slides(self, data: Dict):
        kpis = data.get('kpis', [])
        if not kpis:
            return
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        title.text = "Key Performance Indicators"
        rows = len(kpis) + 1
        cols = 5
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        headers = ['KPI', 'Target', 'Actual', 'Progress', 'Status']
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        for row_idx, kpi in enumerate(kpis[:20], 1):
            table.cell(row_idx, 0).text = str(kpi.get('name', ''))[:30]
            table.cell(row_idx, 1).text = str(kpi.get('target', ''))
            table.cell(row_idx, 2).text = str(kpi.get('actual', ''))
            table.cell(row_idx, 3).text = f"{kpi.get('progress', 0)}%"
            status = kpi.get('status', '')
            cell = table.cell(row_idx, 4)
            cell.text = status
            if status == 'On Track':
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(16, 185, 129)
            elif status == 'At Risk':
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 158, 11)
            elif status == 'Off Track':
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(239, 68, 68)

    def _add_chart_slides(self, data: Dict):
        charts = data.get('charts', [])
        for chart_data in charts[:3]:
            slide = self.prs.slides.add_slide(self.title_only_slide_layout)
            title = slide.shapes.title
            title.text = chart_data.get('title', 'Performance Chart')
            chart_type = chart_data.get('type', 'bar')
            chart_data_obj = chart_data.get('data', {})
            labels = chart_data_obj.get('labels', [])
            values = chart_data_obj.get('values', [])
            if not labels or not values:
                continue
            chart_data_series = CategoryChartData()
            chart_data_series.categories = labels
            chart_data_series.add_series('Value', values)
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
            chart_type_map = {
                'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE
            }
            chart = slide.shapes.add_chart(
                chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
                x, y, cx, cy,
                chart_data_series
            ).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    def _add_table_slides(self, data: Dict):
        tables = data.get('tables', [])
        for table_data in tables[:2]:
            slide = self.prs.slides.add_slide(self.content_slide_layout)
            title = slide.shapes.title
            title.text = table_data.get('title', 'Data Table')
            columns = table_data.get('columns', [])
            rows = table_data.get('rows', [])
            if not columns or not rows:
                continue
            row_count = min(len(rows) + 1, 15)
            col_count = len(columns)
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            table = slide.shapes.add_table(row_count, col_count, left, top, width, height).table
            for col, header in enumerate(columns):
                table.cell(0, col).text = str(header)
                table.cell(0, col).text_frame.paragraphs[0].font.bold = True
                table.cell(0, col).text_frame.paragraphs[0].font.size = Pt(11)
            for row_idx, row in enumerate(rows[:row_count-1], 1):
                for col_idx, value in enumerate(row[:col_count]):
                    table.cell(row_idx, col_idx).text = str(value)

    def _add_mission_status_slides(self, data: Dict):
        mission_data = data.get('mission_status', {})
        if not mission_data:
            return
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        title.text = "Mission Status Report"
        content = slide.placeholders[1]
        mission_text = f"Overall Status: {mission_data.get('status', 'N/A')}\n\n"
        mission_text += f"Performance Analysis:\n{mission_data.get('analysis', 'N/A')}\n\n"
        mission_text += f"Key Challenges:\n{mission_data.get('challenges', 'N/A')}\n\n"
        mission_text += f"Action Plans:\n{mission_data.get('actions', 'N/A')}"
        content.text = mission_text

    def _add_conclusion_slide(self):
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        title.text = "Conclusion"
        content = slide.placeholders[1]
        content.text = "This report was generated by Falcon PMS - Performance Management System.\n\n© Falcon Consulting"

    def export_to_bytes(self, data: Dict[str, Any], report_name: str) -> bytes:
        self.export(data, report_name)
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()